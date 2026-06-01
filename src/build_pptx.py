"""Build a Google-Slides-compatible .pptx mirroring the HTML deck.

Output: slides.pptx in the project root. Upload to Google Drive and open
with Google Slides; it auto-converts.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path("/Users/moin/Robots_diffusion_planner")
OUT = ROOT / "slides.pptx"

W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(0x0C, 0x10, 0x19)
PANEL = RGBColor(0x16, 0x1B, 0x26)
INK = RGBColor(0xE7, 0xEC, 0xF3)
DIM = RGBColor(0x8D, 0x97, 0xA8)
ACCENT = RGBColor(0x6F, 0xB6, 0xFF)
DIFF = RGBColor(0xC5, 0x76, 0xFF)
GOOD = RGBColor(0x7E, 0xE7, 0x87)
BAD = RGBColor(0xFF, 0x7B, 0x72)
BASE = RGBColor(0xFF, 0xB8, 0x6C)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return slide


def add_text(slide, x, y, w, h, text, size=24, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Inter"
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=INK, bullet_color=DIFF):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        rb = p.add_run()
        rb.text = "■  "
        rb.font.size = Pt(size - 2)
        rb.font.color.rgb = bullet_color
        rb.font.bold = True
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Inter"


def add_tag(slide, x, y, text, color=ACCENT, bg_color=None):
    if bg_color is None:
        bg_color = RGBColor(0x2A, 0x3A, 0x55)
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.6), Inches(0.34))
    tag.fill.solid()
    tag.fill.fore_color.rgb = bg_color
    tag.line.fill.background()
    tag.adjustments[0] = 0.5
    tf = tag.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = "Inter"


def add_stat_card(slide, x, y, w, h, big, lbl, big_color=DIFF):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = RGBColor(0x2A, 0x31, 0x42)
    card.adjustments[0] = 0.08
    tf = card.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = big
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = big_color
    r.font.name = "Inter"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = lbl
    r2.font.size = Pt(9)
    r2.font.color.rgb = DIM
    r2.font.name = "Inter"


def add_image(slide, x, y, w, h, path):
    p = ROOT / path
    if not p.exists():
        return None
    pic = slide.shapes.add_picture(str(p), x, y, w, h)
    return pic


def add_pagenum(slide, n, total):
    add_text(
        slide, W - Inches(0.9), H - Inches(0.35), Inches(0.8), Inches(0.25),
        f"{n} / {total}", size=10, color=DIM, align=PP_ALIGN.RIGHT,
    )


def add_corner(slide):
    add_text(
        slide, Inches(0.4), H - Inches(0.35), Inches(6), Inches(0.25),
        "Moin Mattar  --  Diffusion-Guided Frontier Exploration",
        size=10, color=DIM,
    )


TOTAL = 13


def slide_1(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.6), "COSC 81/281 final project")
    add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(2),
             "Diffusion-Guided\nFrontier Exploration",
             size=64, bold=True, color=DIFF)
    add_text(s, Inches(0.6), Inches(3.4), Inches(12), Inches(1.5),
             "Can a generative model of buildings help a robot decide where to\n"
             "drive next, and just as importantly, when can it not?",
             size=22, color=DIM)
    add_text(s, Inches(0.6), Inches(5.8), Inches(12), Inches(0.5),
             "Moin Mattar", size=22, bold=True)
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.3),
             "Dartmouth COSC 81/281  --  Spring 2026  --  June 1",
             size=13, color=DIM)
    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.3),
             "AI helped me in formatting and writing (HTML/CSS), as well as explained concepts.",
             size=10, color=DIM)
    add_pagenum(s, 1, TOTAL)


def slide_2(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "The setup")
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1.8),
             "A robot at the door of a building it has never seen.",
             size=42, bold=True)
    add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(1),
             "2D lidar. After each ping, a slightly bigger map.\n"
             "The question is always the same: which frontier do I drive to next?",
             size=20, color=DIM)
    add_bullets(s, Inches(0.6), Inches(4.4), Inches(12), Inches(2.5), [
        "Classical baseline: score each frontier by unknown cells nearby minus a distance penalty. Memoryless.",
        "What it throws away: walls run straight, hallways connect rooms, doors cluster. Buildings have structure.",
        "Our hypothesis: a model trained on real floor plans can imagine plausible completions and score frontiers under them.",
    ], size=17, bullet_color=BASE)
    add_corner(s); add_pagenum(s, 2, TOTAL)


def slide_3(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "Act 1 -- what I built")
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1.4),
             "The data: HouseExpo, 35k floor plans",
             size=36, bold=True)
    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1),
             "For each plan I generate a partial view from a random lidar scan and pair it with the hidden remainder.\n"
             "Augmentations: 4 rotations, horizontal and vertical flips. Final training set: 2.66M (partial, hidden) pairs.",
             size=15, color=DIM)
    add_image(s, Inches(0.6), Inches(3.6), Inches(12.1), Inches(3.2),
              "results/data_artifacts/01_sample_grid_20.png")
    add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
             "Per pair, left to right: ground truth occupancy grid, partial map (input), hidden remainder (target).",
             size=11, color=DIM)
    add_corner(s); add_pagenum(s, 3, TOTAL)


def slide_4(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "Act 1 -- what I built", color=DIFF,
            bg_color=RGBColor(0x2E, 0x1D, 0x44))
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1),
             "The model: conditional U-Net, 4.16M params",
             size=34, bold=True)
    add_bullets(s, Inches(0.6), Inches(2.4), Inches(7), Inches(4), [
        "Architecture: conditional U-Net with sinusoidal time embedding, channel mults (1, 2, 4, 4), base 32 channels.",
        "Training objective: DDPM with MSE loss on the noise (Ho et al. 2020). T = 1000 diffusion timesteps.",
        "Sampling: DDIM with 30 steps for inference. K = 8 completions per partial map.",
        "Hardware: single T4 GPU, batch size 32, AdamW.",
    ], size=15)
    add_stat_card(s, Inches(8.5), Inches(2.4), Inches(4.2), Inches(1.2),
                  "4.16M", "PARAMETERS IN THE U-NET TRUNK", big_color=DIFF)
    add_stat_card(s, Inches(8.5), Inches(3.8), Inches(4.2), Inches(1.2),
                  "29 epochs", "TRAINED ON 2.66M PAIRS AT 256x256", big_color=ACCENT)
    add_stat_card(s, Inches(8.5), Inches(5.2), Inches(4.2), Inches(1.2),
                  "~2 s", "PER K=8 COMPLETION ON T4 WITH DDIM", big_color=GOOD)
    add_corner(s); add_pagenum(s, 4, TOTAL)


def slide_5(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "Act 1 -- what I built")
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1),
             "Training converged from 0.033 to 0.0007",
             size=32, bold=True)
    add_image(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.7),
              "docs/fig_loss.png")
    add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.8),
             "Left: MSE on noise dropped two orders of magnitude. Right: validation IoU averaged 0.66, best epoch 0.77.\n"
             "The model is learning building structure, not memorising the training set.",
             size=13, color=DIM)
    add_corner(s); add_pagenum(s, 5, TOTAL)


def slide_6(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "Act 1 -- what I built", color=DIFF,
            bg_color=RGBColor(0x2E, 0x1D, 0x44))
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1),
             "Sample quality across epochs: the model wakes up",
             size=30, bold=True)
    add_image(s, Inches(0.4), Inches(2.6), Inches(12.5), Inches(3.4),
              "results/showcase/06_training_progression.png")
    add_text(s, Inches(0.6), Inches(6.15), Inches(12), Inches(0.9),
             "Same partial inputs, different checkpoints. Epoch 5: blurry room-sized blobs. Epoch 10: walls emerge.\n"
             "By epoch 20 the predictions have crisp doorways and corridors. I used epoch 20 for deployment.",
             size=13, color=DIM)
    add_corner(s); add_pagenum(s, 6, TOTAL)


def slide_7(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "Act 1 -- what I built")
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1),
             "K = 8 imagined buildings. Variance is the signal.",
             size=30, bold=True)
    add_image(s, Inches(0.4), Inches(2.4), Inches(12.5), Inches(3.5),
              "results/showcase/03_sample_diversity.png")
    add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.9),
             "Each row is one partial map; columns show eight DDIM samples from the same input. Where the partial map\n"
             "already disambiguates structure, samples agree; where it does not, they diverge. That disagreement is our UCB bonus.",
             size=13, color=DIM)
    add_corner(s); add_pagenum(s, 7, TOTAL)


def slide_8(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "Act 1 -- what I built", color=DIFF,
            bg_color=RGBColor(0x2E, 0x1D, 0x44))
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1),
             "The pipeline: partial map in, argmax frontier out",
             size=30, bold=True)
    add_text(s, Inches(0.6), Inches(2.05), Inches(12), Inches(0.7),
             "Same scoring loop as a classical frontier explorer. The only difference: frontiers are evaluated\n"
             "against K=8 imagined completions, not the raw partial map.",
             size=14, color=DIM)

    box_w = Inches(2.9)
    box_h = Inches(3.0)
    box_y = Inches(3.6)
    gap = Inches(0.25)
    titles = [
        ("STEP 1", "Partial map",
         "Robot's current occupancy grid. Free / unknown / walls. Frontiers on the boundary.", ACCENT),
        ("STEP 2 -- DIFFUSION", "K = 8 completions",
         "Conditional U-Net, 4M params. One batched DDIM pass samples K plausible buildings.", DIFF),
        ("STEP 3 -- SCORE", "E[gain] + UCB - dist",
         "Mean lidar gain across K, plus standard deviation, minus travel distance.", INK),
        ("STEP 4", "Drive argmax",
         "Pick the frontier with the highest score. Move. Sense. Repeat.", ACCENT),
    ]
    x = Inches(0.4)
    for step, label, detail, accent in titles:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL
        box.line.color.rgb = RGBColor(0x2A, 0x31, 0x42)
        box.adjustments[0] = 0.05
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = step
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = accent
        p2 = tf.add_paragraph(); p2.space_before = Pt(8)
        r2 = p2.add_run(); r2.text = label
        r2.font.size = Pt(18); r2.font.bold = True; r2.font.color.rgb = INK
        p3 = tf.add_paragraph(); p3.space_before = Pt(10)
        r3 = p3.add_run(); r3.text = detail
        r3.font.size = Pt(11); r3.font.color.rgb = DIM
        x += box_w + gap
    add_corner(s); add_pagenum(s, 8, TOTAL)


def slide_9(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "Act 2 -- what the prof asked for",
            color=BAD, bg_color=RGBColor(0x3A, 0x1C, 0x20))
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1.1),
             "What I showed him -- and the four asks that followed",
             size=28, bold=True)

    add_text(s, Inches(0.6), Inches(2.4), Inches(6), Inches(0.4),
             "What I showed at the check-in", size=18, bold=True, color=INK)
    add_bullets(s, Inches(0.6), Inches(3.0), Inches(6), Inches(4), [
        "The data, the U-Net, the loss curve, K=8 sampling.",
        "An early result on one map: 'baseline is good. Diffusion is maybe better at the beginning. But the baseline mostly catches up.'",
        "Honest framing: the prior helps in the early-budget regime; it does not beat the heuristic asymptotically.",
    ], size=14, bullet_color=BASE)

    add_text(s, Inches(7.0), Inches(2.4), Inches(6), Inches(0.4),
             "His four asks", size=18, bold=True, color=DIFF)
    add_bullets(s, Inches(7.0), Inches(3.0), Inches(6), Inches(4), [
        "1. Compare against four scorers: nearest, info-gain only, info-gain + distance, ours.",
        "2. Show it running in a real simulator. Stage (used in PA3) is acceptable.",
        "3. Report actual inference timing.",
        "4. Describe limitations honestly in the report.",
    ], size=14, bullet_color=DIFF)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             "The remaining slides walk through how I addressed each ask, plus the deeper experiments I added on top.",
             size=14, color=DIM)
    add_corner(s); add_pagenum(s, 9, TOTAL)


def slide_10(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "His ask #1 -- 4-baseline ablation",
            color=DIFF, bg_color=RGBColor(0x2E, 0x1D, 0x44))
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1.0),
             "Four scorers, same 30 maps. Clean decomposition.",
             size=28, bold=True)
    add_image(s, Inches(0.4), Inches(2.3), Inches(7.2), Inches(4.2),
              "results/analysis/figures/12_4baseline_curves.png")
    add_text(s, Inches(0.4), Inches(6.6), Inches(7.2), Inches(0.4),
             "N=30 HouseExpo, lidar=70, identical seeds. Bands = mean +/- 1 SE.",
             size=10, color=DIM)
    add_stat_card(s, Inches(7.9), Inches(2.3), Inches(4.9), Inches(1.25),
                  "+3.5pp", "DIFFUSION VS HEURISTIC AT STEP 4. THE LEARNED PRIOR ON TOP.", big_color=GOOD)
    add_stat_card(s, Inches(7.9), Inches(3.7), Inches(4.9), Inches(1.25),
                  "+17.5pp", "INFO-GAIN VS NEAREST. INFO-GAIN IS DOING MOST OF THE WORK.", big_color=ACCENT)
    add_stat_card(s, Inches(7.9), Inches(5.1), Inches(4.9), Inches(1.25),
                  "~0pp", "DISTANCE TERM ON TOP OF INFO-GAIN. DECORATIVE ON HOUSEEXPO.", big_color=INK)
    add_corner(s); add_pagenum(s, 10, TOTAL)


def slide_11(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "His ask #2 -- show it running",
            color=DIFF, bg_color=RGBColor(0x2E, 0x1D, 0x44))
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1.0),
             "World -> observation -> imagination -> pick",
             size=28, bold=True)
    add_image(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(4.3),
              "results/analysis/figures/15_behind_mind_map2638.png")
    add_text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.7),
             "Top: world + robot trail. Middle: partial map + frontier candidates. "
             "Bottom: mean of K=4 diffusion samples (the imagined building) with the chosen frontier circled.\n"
             "On map 2638 the prior carries the robot to 95% coverage at step 4.",
             size=12, color=DIM)
    add_corner(s); add_pagenum(s, 11, TOTAL)


def slide_12(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "His ask #4 -- limits + timing",
            color=BAD, bg_color=RGBColor(0x3A, 0x1C, 0x20))
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1.1),
             "When the prior fails, and what realtime would buy",
             size=26, bold=True)
    add_image(s, Inches(0.4), Inches(2.2), Inches(6.3), Inches(3.6),
              "results/analysis/figures/07_in_vs_ood_delta.png")
    add_text(s, Inches(0.4), Inches(5.9), Inches(6.3), Inches(1.1),
             "OOD ablation. A second U-Net trained on synthetic warehouse layouts, "
             "evaluated on the same 30 HouseExpo maps. Step-4 advantage drops from +3.5pp to -0.1pp. "
             "The prior is not a generic averaging trick.",
             size=10, color=DIM)
    add_image(s, Inches(7.0), Inches(2.2), Inches(6.0), Inches(3.6),
              "results/analysis/figures/11_2x2_matrix.png")
    add_text(s, Inches(7.0), Inches(5.9), Inches(6.0), Inches(1.1),
             "Realtime upper bound. Re-sampling K=8 every 3 driving cells "
             "(what a distilled model could achieve) lifts in-domain advantage to +10.0pp. "
             "Prior contributes ~4pp; sampling frequency contributes ~7pp.",
             size=10, color=DIM)
    add_text(s, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4),
             "Inference timing on T4: ~2s for K=8 at DDIM=30. K=4 + DDIM=10 measured at ~500ms -- exactly the realtime budget.",
             size=12, bold=True, color=DIFF)
    add_corner(s); add_pagenum(s, 12, TOTAL)


def slide_13(prs):
    s = new_slide(prs)
    add_tag(s, Inches(0.6), Inches(0.5), "Conclusion")
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1),
             "Where this is useful, and what is next",
             size=32, bold=True)

    add_text(s, Inches(0.6), Inches(2.3), Inches(6), Inches(0.4),
             "Useful when all three hold", size=18, bold=True, color=GOOD)
    add_bullets(s, Inches(0.6), Inches(2.9), Inches(6), Inches(3.5), [
        "Deployment environment is structured and known in advance (warehouses, hospitals, office buildings).",
        "Robot has a tight step budget (battery, time, danger window). The advantage concentrates in steps 2-5.",
        "Inference is fast enough to drive (sub-second). Latency measurements show this is achievable on T4.",
    ], size=13, bullet_color=GOOD)

    add_text(s, Inches(7.0), Inches(2.3), Inches(6), Inches(0.4),
             "Honest limits", size=18, bold=True, color=BAD)
    add_bullets(s, Inches(7.0), Inches(2.9), Inches(6), Inches(3.5), [
        "Asymptotic coverage. The baseline catches up by step 20.",
        "OOD priors are useless. Training distribution must match deployment.",
        "2/30 hard failures. K-sample disagreement could flag these and fall back to baseline.",
        "Live hardware deployment. Stage integration uses pre-computed waypoints; Docker QEMU is 50x too slow for live diffusion.",
    ], size=13, bullet_color=BAD)

    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.6),
             "Thank you. Questions?", size=26, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
             "Moin Mattar  --  AI helped me in formatting and writing (HTML/CSS) and explained concepts.",
             size=10, color=DIM, align=PP_ALIGN.CENTER)
    add_pagenum(s, 13, TOTAL)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for fn in [slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7,
               slide_8, slide_9, slide_10, slide_11, slide_12, slide_13]:
        fn(prs)

    prs.save(str(OUT))
    print(f"Wrote {OUT} ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
