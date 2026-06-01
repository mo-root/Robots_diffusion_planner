"""Journey-style PPTX matching the user's billiards presentation rhythm.

White background, big section dividers ("1 | Novelty & Motivation"),
Project Goal / Novelty / Hypothesis / Pipeline 4-cards / Math objectives /
Now you tell me / Hardships / Sampling timing / DEMO / What's Next / Thanks.

Output: journey_deck.pptx — upload to Google Drive, open with Google Slides.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path("/Users/moin/Robots_diffusion_planner")
OUT = ROOT / "journey_deck.pptx"

W = Inches(13.333)
H = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DIM = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x1F, 0x6B, 0x47)
ORANGE = RGBColor(0xD9, 0x76, 0x06)
BLUE = RGBColor(0x25, 0x63, 0xEB)
YELLOW = RGBColor(0xD4, 0xA0, 0x17)
RED = RGBColor(0xB9, 0x1C, 0x1C)
LIGHT = RGBColor(0xF4, 0xF7, 0xF3)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    return s


def text(slide, x, y, w, h, content, size=18, bold=False, italic=False,
         color=BLACK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def title_bar(slide, title_text, color=BLACK):
    text(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.7),
         title_text, size=36, bold=True, color=color)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(1.18), Inches(12.1), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()


def bullets(slide, x, y, w, h, items, size=16, color=BLACK, bullet_color=GREEN):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)

        parts = []
        rest = item
        while "**" in rest:
            before, _, rest2 = rest.partition("**")
            bold_text, _, rest = rest2.partition("**")
            if before:
                parts.append((before, False))
            parts.append((bold_text, True))
        if rest:
            parts.append((rest, False))

        br = p.add_run()
        br.text = "•  "
        br.font.size = Pt(size); br.font.bold = True
        br.font.color.rgb = bullet_color; br.font.name = "Calibri"

        for t, bold in parts:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"


def picture(slide, x, y, w, h, path):
    p = ROOT / path
    if not p.exists():
        print(f"missing: {p}")
        return None
    return slide.shapes.add_picture(str(p), x, y, w, h)


def pagenum(slide, n):
    text(slide, W - Inches(0.6), H - Inches(0.4), Inches(0.4), Inches(0.3),
         str(n), size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
def slide_1(prs):
    s = new_slide(prs)
    text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.35),
         "COSC 81/281  ·  FINAL PROJECT  ·  SPRING 2026",
         size=11, bold=True, color=DIM)
    text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(2.2),
         "Diffusion-Guided\nFrontier Exploration",
         size=56, bold=True, color=BLACK)
    # Green underbar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), Inches(4.45), Inches(2.2), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(1.0),
         "Can a generative model of buildings help a robot decide where to drive\n"
         "next — and just as importantly, when can it not?",
         size=20, italic=True, color=DIM)
    text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
         "Moin Mattar", size=20, bold=True, color=BLACK)
    text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.3),
         "Dartmouth  ·  June 1, 2026", size=12, color=DIM)
    text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.3),
         "AI helped me in formatting and writing (HTML/CSS/LaTeX), as well as explained concepts.",
         size=9, color=DIM)
    pagenum(s, 1)


# ============================================================
# SLIDE 2 — SECTION DIVIDER "1 | Novelty & Motivation"
# ============================================================
def section_divider(prs, n, title, sub):
    s = new_slide(prs)
    text(s, Inches(2.5), Inches(2.7), Inches(1.5), Inches(2),
         str(n), size=120, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    # vertical bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(4.0), Inches(2.5), Inches(0.06), Inches(2.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()
    text(s, Inches(4.4), Inches(2.7), Inches(8.5), Inches(1.5),
         title, size=64, bold=True, color=GREEN)
    text(s, Inches(4.4), Inches(4.2), Inches(8.5), Inches(0.6),
         sub, size=20, italic=True, color=DIM)
    return s


def slide_2(prs):
    s = section_divider(prs, 1, "Novelty & Motivation",
                        "A robot at the door of a building it has never seen.")
    pagenum(s, 2)


# ============================================================
# SLIDE 3 — PROJECT GOAL
# ============================================================
def slide_3(prs):
    s = new_slide(prs)
    title_bar(s, "Project Goal")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.4),
         "Autonomous mapping of an unknown building, one frontier decision at a time.",
         size=14, italic=True, color=DIM)

    items = [
        "**Setup:** a robot is dropped at the door of a building it has never seen, with only a 2D lidar.",
        "**Input each step:** the partial occupancy grid the robot has accumulated.",
        "**Output each step:** the next frontier to drive to.",
        "**Assumption:** deployment environment shares structural distribution with training data.",
        "**What's broken:** classical frontier scoring is memoryless — it throws away the fact that real buildings have walls that run straight, hallways that connect rooms, doors that cluster.",
    ]
    bullets(s, Inches(0.6), Inches(2.1), Inches(7.5), Inches(5), items, size=14)
    picture(s, Inches(8.4), Inches(2.0), Inches(4.4), Inches(4.8),
            "results/analysis/figures/15_behind_mind_map2638.png")
    pagenum(s, 3)


# ============================================================
# SLIDE 4 — NOVELTY?
# ============================================================
def slide_4(prs):
    s = new_slide(prs)
    title_bar(s, "Novelty?")
    items = [
        "We treat **diffusion samples as imagined buildings**, not just predictions. K=8 plausible completions per frontier decision.",
        "We use **K-sample disagreement as an upper-confidence-bound exploration signal** — the variance is the feature, not noise to average out.",
        "We isolate **where the wins come from** with a 4-baseline ablation, plus an **out-of-domain control** that proves the prior is doing real structural work.",
    ]
    bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(4), items, size=18)
    pagenum(s, 4)


# ============================================================
# SLIDE 5 — HYPOTHESIS (boxed)
# ============================================================
def slide_5(prs):
    s = new_slide(prs)
    title_bar(s, "Hypothesis")

    # green left border
    border = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.6), Inches(2.8), Inches(0.08), Inches(2.6))
    border.fill.solid(); border.fill.fore_color.rgb = GREEN
    border.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.95), Inches(2.85), Inches(8), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.line_spacing = 1.4

    r1 = p.add_run()
    r1.text = "A learned generative prior of building layouts gives a robot a measurable early-budget head start in frontier exploration "
    r1.font.size = Pt(20); r1.font.color.rgb = BLACK; r1.font.name = "Calibri"

    r2 = p.add_run()
    r2.text = "— but only when "
    r2.font.size = Pt(20); r2.font.italic = True; r2.font.color.rgb = BLACK; r2.font.name = "Calibri"

    r3 = p.add_run()
    r3.text = "(1) the training distribution matches deployment, (2) the step budget is tight, and (3) inference is fast enough to drive."
    r3.font.size = Pt(20); r3.font.bold = True; r3.font.color.rgb = BLACK; r3.font.name = "Calibri"

    # right image — small diversity sample
    picture(s, Inches(9.3), Inches(2.6), Inches(3.6), Inches(3.4),
            "results/showcase/03_sample_diversity.png")
    pagenum(s, 5)


# ============================================================
# SLIDE 6 — SECTION DIVIDER "2 | Our System"
# ============================================================
def slide_6(prs):
    s = section_divider(prs, 2, "Our System",
                        "HouseExpo  →  Conditional U-Net (DDPM)  →  K-sample scoring  →  Frontier pick")
    pagenum(s, 6)


# ============================================================
# SLIDE 7 — END-TO-END PIPELINE (4 cards like billiards slide 9)
# ============================================================
def slide_7(prs):
    s = new_slide(prs)
    title_bar(s, "End-to-End Pipeline I Implemented")

    cards = [
        ("01", "Data\nCollection",
         "HouseExpo: 35,126 residential floor plans. Per plan, partial views + hidden remainders → 2.66M training pairs after augmentation.", GREEN),
        ("02", "U-Net\n+ DDPM",
         "Conditional U-Net, 4.16M params. MSE on noise. 29 epochs on T4 GPU. Loss 0.033 → 0.0007.", ORANGE),
        ("03", "K-Sample\nScoring",
         "K=8 DDIM completions. Score = E[gain] + λ·Std[gain] − β·dist. Variance = UCB exploration bonus.", YELLOW),
        ("04", "ROS 2\nIntegration",
         "ROS 2 node publishes /best_frontier. Exploration manager (from PA3) drives robot in Stage maze.", BLUE),
    ]

    card_w = Inches(2.95)
    card_h = Inches(4.5)
    x = Inches(0.7)
    y = Inches(2.0)
    gap = Inches(0.15)

    for num, title, body, accent in cards:
        # main card
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = accent
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.05

        # top accent tab
        tab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.1))
        tab.fill.solid(); tab.fill.fore_color.rgb = accent
        tab.line.fill.background()

        # number
        text(s, x + Inches(0.25), y + Inches(0.3), card_w - Inches(0.5), Inches(0.9),
             num, size=44, bold=True, color=accent)
        # title
        text(s, x + Inches(0.25), y + Inches(1.3), card_w - Inches(0.5), Inches(1.2),
             title, size=22, bold=True, color=BLACK)
        # body
        text(s, x + Inches(0.25), y + Inches(2.7), card_w - Inches(0.5), card_h - Inches(2.8),
             body, size=11, color=DIM)

        x += card_w + gap

    pagenum(s, 7)


# ============================================================
# SLIDE 8 — DATA (table + sample image)
# ============================================================
def slide_8(prs):
    s = new_slide(prs)
    title_bar(s, "The Data — HouseExpo")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.4),
         "Augmentation (4 rotations + flips) was the dominant lever for sample diversity.",
         size=14, italic=True, color=DIM)

    # Simple table
    rows = [
        ["Source", "Plans", "Pairs after augmentation", "Purpose"],
        ["HouseExpo (residential)", "35,126", "2.66M", "training"],
        ["Held-out HouseExpo", "30", "—", "evaluation"],
        ["Synthetic warehouse", "—", "small set", "OOD control"],
    ]
    table_x = Inches(0.6)
    table_y = Inches(2.1)
    col_w = [Inches(3.2), Inches(1.5), Inches(2.6), Inches(1.7)]
    row_h = Inches(0.5)

    for ri, row in enumerate(rows):
        cx = table_x
        is_header = ri == 0
        for ci, cell in enumerate(row):
            if is_header:
                bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, table_y + ri * row_h, col_w[ci], row_h)
                bg.fill.solid(); bg.fill.fore_color.rgb = GREEN
                bg.line.color.rgb = WHITE
            else:
                line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, table_y + ri * row_h + row_h - Inches(0.02), col_w[ci], Inches(0.01))
                line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                line.line.fill.background()
            tb = s.shapes.add_textbox(cx + Inches(0.1), table_y + ri * row_h + Inches(0.1),
                                      col_w[ci] - Inches(0.2), row_h - Inches(0.2))
            tf = tb.text_frame
            tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = cell
            r.font.size = Pt(13)
            r.font.bold = is_header
            r.font.color.rgb = WHITE if is_header else BLACK
            r.font.name = "Calibri"
            cx += col_w[ci]

    # Sample image on right
    picture(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.6),
            "results/data_artifacts/01_sample_grid_20.png")

    pagenum(s, 8)


# ============================================================
# SLIDE 9 — BASELINE OBJECTIVE (DDPM math, like billiards slide 15)
# ============================================================
def slide_9(prs):
    s = new_slide(prs)
    title_bar(s, "Baseline Objective (DDPM)")

    text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.5),
         "Forward process:", size=16, color=BLACK)
    # math line 1
    text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.6),
         "x_τ = √(α̅_τ) · x₀ + √(1 − α̅_τ) · ε,    ε ~ N(0, I)",
         size=20, italic=True, color=BLACK, align=PP_ALIGN.CENTER)

    text(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.5),
         "One-step objective:", size=16, color=BLACK)
    text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.6),
         "L_base = E [ ‖ε̂_θ(x_τ, τ, partial)  −  ε‖² ]",
         size=20, italic=True, color=BLACK, align=PP_ALIGN.CENTER)

    items = [
        "Conditioned on the **partial occupancy grid** + a **known-mask** channel.",
        "Single-step denoising target — MSE on the noise added at step τ.",
        "DDIM at inference: **30 steps**, K=8 samples per batched call (~2 s on T4).",
    ]
    bullets(s, Inches(0.6), Inches(5.2), Inches(12), Inches(2), items, size=14)
    pagenum(s, 9)


# ============================================================
# SLIDE 10 — OUR SCORING OBJECTIVE (K-sample UCB)
# ============================================================
def slide_10(prs):
    s = new_slide(prs)
    title_bar(s, "Frontier Scoring Objective")

    text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.5),
         "For each frontier cluster f, with K=8 imagined completions:", size=16, color=BLACK)

    text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.7),
         "score(f)  =  E_k[ G_k(f) ]  +  λ · Std_k[ G_k(f) ]  −  β · dist(robot, f)",
         size=18, italic=True, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    text(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.6),
         "where G_k(f) = unknown cells inside frontier sensor footprint under completion k.",
         size=13, color=DIM, align=PP_ALIGN.CENTER)

    items = [
        "**Mean gain** — how much the robot expects to discover at this frontier.",
        "**Variance bonus** — UCB-style; the model's disagreement is the exploration signal.",
        "**Distance term** — travel cost.",
        "Pick the argmax. Move. Sense. Re-sample. Repeat.",
    ]
    bullets(s, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5), items, size=15)
    pagenum(s, 10)


# ============================================================
# SLIDE 11 — TRAINING & SAMPLES (the journey of training)
# ============================================================
def slide_11(prs):
    s = new_slide(prs)
    title_bar(s, "Training Journey:  loss  0.033 → 0.0007")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.4),
         "29 epochs on T4. Validation IoU averaged 0.66, best 0.77.",
         size=14, italic=True, color=DIM)

    picture(s, Inches(0.5), Inches(2.0), Inches(7.0), Inches(3.0),
            "docs/fig_loss.png")
    picture(s, Inches(7.7), Inches(2.0), Inches(5.2), Inches(2.0),
            "results/showcase/06_training_progression.png")
    text(s, Inches(7.7), Inches(4.1), Inches(5.2), Inches(0.4),
         "Sample quality at epochs 5/10/15/20/25 — the model wakes up.",
         size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)

    # K=8 diversity below
    picture(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7),
            "results/showcase/03_sample_diversity.png")
    text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
         "K=8 DDIM samples per partial input — variance becomes the exploration signal.",
         size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)

    pagenum(s, 11)


# ============================================================
# SLIDE 12 — "NOW YOU TELL ME" 4-BASELINE STRIP
# ============================================================
def slide_12(prs):
    s = new_slide(prs)
    title_bar(s, "Now you tell me — where do the wins come from?")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.4),
         "Four scorers, 30 HouseExpo maps, identical seeds.",
         size=14, italic=True, color=DIM)

    picture(s, Inches(0.5), Inches(2.0), Inches(8.5), Inches(4.6),
            "results/analysis/figures/12_4baseline_curves.png")

    # Stat cards on the right
    cards = [
        ("+3.5pp", "DIFFUSION VS HEURISTIC AT STEP 4.\nThe learned prior on top of info-gain + distance.", GREEN),
        ("+17.5pp", "INFO-GAIN VS NEAREST AT STEP 4.\nInfo-gain is doing most of the work.", BLUE),
        ("~0pp", "DISTANCE TERM CONTRIBUTION.\nDecorative on HouseExpo.", ORANGE),
    ]
    cy = Inches(2.0)
    for big, lbl, accent in cards:
        # left border
        border = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(9.2), cy, Inches(0.08), Inches(1.4))
        border.fill.solid(); border.fill.fore_color.rgb = accent
        border.line.fill.background()
        text(s, Inches(9.4), cy + Inches(0.05), Inches(3.5), Inches(0.6),
             big, size=28, bold=True, color=accent)
        text(s, Inches(9.4), cy + Inches(0.7), Inches(3.5), Inches(0.7),
             lbl, size=9, color=DIM)
        cy += Inches(1.55)

    pagenum(s, 12)


# ============================================================
# SLIDE 13 — RESULTS: BEHIND THE MIND
# ============================================================
def slide_13(prs):
    s = new_slide(prs)
    title_bar(s, "Behind the mind — the loop made visible")

    picture(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.0),
            "results/analysis/figures/15_behind_mind_map2638.png")
    text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
         "Map 2638 — top: world + trail. Middle: partial map + frontiers. Bottom: K=4 imagined buildings with picked frontier circled. By step 4 the prior carries the robot to 95% coverage.",
         size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)
    pagenum(s, 13)


# ============================================================
# SLIDE 14 — WHEN THE PRIOR FAILS (OOD + realtime)
# ============================================================
def slide_14(prs):
    s = new_slide(prs)
    title_bar(s, "When the prior fails — and what realtime would buy")

    picture(s, Inches(0.4), Inches(1.8), Inches(6.2), Inches(3.6),
            "results/analysis/figures/07_in_vs_ood_delta.png")
    text(s, Inches(0.4), Inches(5.5), Inches(6.2), Inches(1.2),
         "OOD ablation. Second U-Net trained on synthetic warehouses, evaluated on residential maps. "
         "Step-4 advantage drops from +3.5pp to −0.1pp. The prior is not a generic averaging trick.",
         size=11, italic=True, color=DIM)

    picture(s, Inches(7.0), Inches(1.8), Inches(5.9), Inches(3.6),
            "results/analysis/figures/11_2x2_matrix.png")
    text(s, Inches(7.0), Inches(5.5), Inches(5.9), Inches(1.2),
         "Realtime upper bound. Re-sampling K=8 every 3 driving cells (the speed a distilled model could achieve) "
         "lifts in-domain advantage to +10.0pp. Prior contributes ~4pp; sampling frequency contributes ~7pp.",
         size=11, italic=True, color=DIM)

    pagenum(s, 14)


# ============================================================
# SLIDE 15 — HARDSHIPS (bulleted, like billiards slide 22)
# ============================================================
def slide_15(prs):
    s = new_slide(prs)
    title_bar(s, "Hardships")

    text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.4),
         "Compute / infra:", size=18, bold=True, color=BLACK)
    items_infra = [
        "**Docker QEMU on Apple Silicon** too slow for live diffusion in Stage — switched to pre-computed waypoints.",
        "**AWS EC2 g4dn.xlarge** (T4 GPU) for training and inference experiments.",
        "**Batched K=8 sampling rewrite** — 3× speedup with no quality loss.",
    ]
    bullets(s, Inches(0.8), Inches(2.5), Inches(12), Inches(2), items_infra, size=14)

    text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.4),
         "Honest limits:", size=18, bold=True, color=RED)
    items_limits = [
        "**Asymptotic coverage** — baseline catches up by step 20.",
        "**Out-of-domain priors** give −0.1pp (the kill experiment).",
        "**2 / 30 hard failures** — K-sample disagreement could flag these.",
        "**Live hardware** — Stage works via pre-computed waypoints, not live inference.",
    ]
    bullets(s, Inches(0.8), Inches(5.2), Inches(12), Inches(2.2), items_limits, size=14, bullet_color=RED)

    pagenum(s, 15)


# ============================================================
# SLIDE 16 — SAMPLING (DDPM/DDIM timing, like billiards slide 23)
# ============================================================
def slide_16(prs):
    s = new_slide(prs)
    title_bar(s, "Sampling")

    items = [
        "**DDPM (1000 steps), K=8:**     ~40 seconds on T4 — way too slow.",
        "**DDIM (30 steps), K=8:**         ~2 seconds (20× faster) — what I used for experiments.",
        "**DDIM (10 steps), K=4:**         ~500 ms (80× faster) — **realtime budget**.",
        "Distillation to Jetson would put **+10pp in-domain advantage** in reach for a real robot.",
    ]
    bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(4), items, size=17)

    # latency profile image on right
    picture(s, Inches(0.6), Inches(5.0), Inches(12.0), Inches(2.0),
            "results/analysis/figures/06b_latency_profile.png")
    pagenum(s, 16)


# ============================================================
# SLIDE 17 — DEMO (side-by-side gif, like billiards slide 24 "DEMO?")
# ============================================================
def slide_17(prs):
    s = new_slide(prs)
    title_bar(s, "DEMO")

    picture(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.0),
            "results/analysis/figures/17_side_by_side_demo.gif")
    text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
         "Top: diffusion-guided.  Bottom: baseline heuristic.  Same map, same start, same lidar. "
         "Both reach high coverage — diffusion is consistently ahead in the early steps.",
         size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)

    pagenum(s, 17)


# ============================================================
# SLIDE 18 — WHAT'S NEXT (like billiards slide 26)
# ============================================================
def slide_18(prs):
    s = new_slide(prs)
    title_bar(s, "What's Next?")
    items = [
        "**Distillation** of the U-Net to Jetson for sub-second inference on a ROSbot.",
        "**Richer training domains** — Gibson, Matterport, real BIM data — for cross-environment generalisation.",
        "**K-sample disagreement** as an \"I don't know\" flag — fall back to heuristic on hard maps.",
        "**Diffusion Forcing for maps** — sequential prediction over multiple exploration timesteps.",
        "**Live ROSbot deployment** in a matched-distribution building (residential / hospital / warehouse).",
    ]
    bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(5), items, size=17)
    pagenum(s, 18)


# ============================================================
# SLIDE 19 — THANKS (meme like billiards slide 28)
# ============================================================
def slide_19(prs):
    s = new_slide(prs)
    text(s, Inches(0.6), Inches(2.8), Inches(12), Inches(1.2),
         "Thanks for your attention.",
         size=54, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.8),
         "Questions?",
         size=36, italic=True, color=GREEN, align=PP_ALIGN.CENTER)

    text(s, Inches(0.6), Inches(5.8), Inches(12), Inches(0.4),
         "Moin Mattar  ·  COSC 81/281  ·  June 1, 2026",
         size=14, color=DIM, align=PP_ALIGN.CENTER)
    text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
         "AI helped me in formatting and writing (HTML/CSS/LaTeX), as well as explained concepts.",
         size=10, color=DIM, align=PP_ALIGN.CENTER)
    pagenum(s, 19)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for fn in [slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7,
               slide_8, slide_9, slide_10, slide_11, slide_12, slide_13,
               slide_14, slide_15, slide_16, slide_17, slide_18, slide_19]:
        fn(prs)

    prs.save(str(OUT))
    print(f"Wrote {OUT} ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
