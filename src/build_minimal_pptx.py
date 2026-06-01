"""Minimal 10-slide PPTX. Short. Clear. Done."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path("/Users/moin/Robots_diffusion_planner")
OUT = ROOT / "minimal_deck.pptx"

W = Inches(13.333); H = Inches(7.5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DIM = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x1F, 0x6B, 0x47)
ORANGE = RGBColor(0xD9, 0x76, 0x06)
BLUE = RGBColor(0x25, 0x63, 0xEB)
YELLOW = RGBColor(0xD4, 0xA0, 0x17)
RED = RGBColor(0xB9, 0x1C, 0x1C)
LIGHT = RGBColor(0xF4, 0xF7, 0xF3)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    return s


def text(slide, x, y, w, h, content, size=18, bold=False, italic=False,
         color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = content
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"


def title_bar(slide, title_text):
    text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7),
         title_text, size=34, bold=True, color=BLACK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()


def bullets(slide, x, y, w, h, items, size=16, color=BLACK, bullet_color=GREEN):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(10)
        parts = []; rest = item
        while "**" in rest:
            before, _, rest2 = rest.partition("**")
            bold_text, _, rest = rest2.partition("**")
            if before: parts.append((before, False))
            parts.append((bold_text, True))
        if rest: parts.append((rest, False))
        br = p.add_run(); br.text = "•  "
        br.font.size = Pt(size); br.font.bold = True
        br.font.color.rgb = bullet_color; br.font.name = "Calibri"
        for t, bold in parts:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "Calibri"


def picture(slide, x, y, w, h, path):
    p = ROOT / path
    if not p.exists():
        print(f"missing: {p}"); return None
    return slide.shapes.add_picture(str(p), x, y, w, h)


def pagenum(slide, n, total=10):
    text(slide, W - Inches(0.7), H - Inches(0.4), Inches(0.5), Inches(0.3),
         f"{n} / {total}", size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ============================================================
# 1 — TITLE
# ============================================================
def s1(prs):
    s = new_slide(prs)
    text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.35),
         "COSC 81/281  ·  FINAL PROJECT", size=11, bold=True, color=DIM)
    text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(2.2),
         "Diffusion-Guided\nFrontier Exploration", size=56, bold=True, color=BLACK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), Inches(4.6), Inches(2.2), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
    text(s, Inches(0.6), Inches(4.85), Inches(12), Inches(1.0),
         "A learned prior of building layouts as a robot's exploration signal.",
         size=22, italic=True, color=DIM)
    text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
         "Moin Mattar  ·  June 1, 2026", size=18, bold=True, color=BLACK)
    text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
         "AI helped me in formatting and writing (HTML/CSS/LaTeX), as well as explained concepts.",
         size=9, color=DIM)
    pagenum(s, 1)


# ============================================================
# 2 — PROJECT GOAL + HYPOTHESIS (combined)
# ============================================================
def s2(prs):
    s = new_slide(prs)
    title_bar(s, "Project Goal  +  Hypothesis")

    text(s, Inches(0.6), Inches(1.4), Inches(7.5), Inches(0.4),
         "The problem", size=16, bold=True, color=GREEN)
    bullets(s, Inches(0.6), Inches(1.85), Inches(7.5), Inches(2.2), [
        "Robot dropped at the door of an unknown building, only 2D lidar.",
        "Input: partial occupancy grid.  Output: next frontier to drive to.",
        "Classical scorer is **memoryless** — throws away building structure.",
    ], size=14)

    text(s, Inches(0.6), Inches(4.3), Inches(7.5), Inches(0.4),
         "Hypothesis", size=16, bold=True, color=GREEN)
    border = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.6), Inches(4.8), Inches(0.08), Inches(2.0))
    border.fill.solid(); border.fill.fore_color.rgb = GREEN; border.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.85), Inches(4.85), Inches(7.4), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.line_spacing = 1.35
    r1 = p.add_run(); r1.text = "A learned generative prior helps "
    r1.font.size = Pt(15); r1.font.color.rgb = BLACK; r1.font.name = "Calibri"
    r2 = p.add_run(); r2.text = "only when "
    r2.font.size = Pt(15); r2.font.italic = True; r2.font.color.rgb = BLACK; r2.font.name = "Calibri"
    r3 = p.add_run(); r3.text = "(1) training distribution matches deployment, (2) the step budget is tight, and (3) inference is fast enough to drive."
    r3.font.size = Pt(15); r3.font.bold = True; r3.font.color.rgb = BLACK; r3.font.name = "Calibri"

    picture(s, Inches(8.6), Inches(1.5), Inches(4.3), Inches(5.4),
            "results/analysis/figures/15_behind_mind_map2638.png")
    pagenum(s, 2)


# ============================================================
# 3 — APPROACH (4-card pipeline)
# ============================================================
def s3(prs):
    s = new_slide(prs)
    title_bar(s, "Approach — End-to-End Pipeline")

    cards = [
        ("01", "Data",
         "HouseExpo: 35k floor plans → 2.66M (partial, hidden) pairs after augmentation.", GREEN),
        ("02", "Model",
         "Conditional U-Net, 4.16M params. DDPM, MSE on noise. 29 epochs on T4 GPU.", ORANGE),
        ("03", "Scoring",
         "K=8 DDIM completions. score = E[gain] + λ·Std[gain] − β·dist.", YELLOW),
        ("04", "Integration",
         "ROS 2 node publishes /best_frontier. Stage maze via pre-computed waypoints.", BLUE),
    ]
    card_w = Inches(2.95); card_h = Inches(4.6)
    x = Inches(0.7); y = Inches(1.9); gap = Inches(0.15)
    for num, title, body, accent in cards:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = accent; card.line.width = Pt(1.5)
        card.adjustments[0] = 0.05
        tab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.1))
        tab.fill.solid(); tab.fill.fore_color.rgb = accent; tab.line.fill.background()
        text(s, x + Inches(0.25), y + Inches(0.35), card_w - Inches(0.5), Inches(0.9),
             num, size=44, bold=True, color=accent)
        text(s, x + Inches(0.25), y + Inches(1.4), card_w - Inches(0.5), Inches(0.7),
             title, size=24, bold=True, color=BLACK)
        text(s, x + Inches(0.25), y + Inches(2.3), card_w - Inches(0.5), card_h - Inches(2.4),
             body, size=12, color=DIM)
        x += card_w + gap
    pagenum(s, 3)


# ============================================================
# 4 — THE LOOP RUNNING (main demo GIF)
# ============================================================
def s4(prs):
    s = new_slide(prs)
    title_bar(s, "DEMO — the loop running")
    picture(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.8),
            "results/exploration_demo/exploration_diffusion.gif")
    text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8),
         "Partial map (with frontiers) → diffusion's prediction → ground truth (robot doesn't see) → coverage growing.",
         size=14, italic=True, color=DIM, align=PP_ALIGN.CENTER)
    pagenum(s, 4)


# ============================================================
# 5 — SIDE-BY-SIDE DEMO
# ============================================================
def s5(prs):
    s = new_slide(prs)
    title_bar(s, "DEMO — diffusion vs baseline, same map")
    picture(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.8),
            "results/analysis/figures/17_side_by_side_demo.gif")
    text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8),
         "Top: diffusion-guided.  Bottom: heuristic baseline. Same map, same start. Diffusion commits to the unseen wing one or two steps earlier.",
         size=14, italic=True, color=DIM, align=PP_ALIGN.CENTER)
    pagenum(s, 5)


# ============================================================
# 6 — RESULTS: 4-BASELINE
# ============================================================
def s6(prs):
    s = new_slide(prs)
    title_bar(s, "Results — where do the wins come from?")
    picture(s, Inches(0.5), Inches(1.6), Inches(8.0), Inches(5.0),
            "results/analysis/figures/12_4baseline_curves.png")
    cards = [
        ("+3.5pp", "Diffusion vs heuristic at step 4.\nThe learned prior on top.", GREEN),
        ("+17.5pp", "Info-gain vs nearest.\nInfo-gain is doing most of the work.", BLUE),
        ("~0pp", "Distance term contribution.\nDecorative on HouseExpo.", ORANGE),
    ]
    cy = Inches(1.8)
    for big, lbl, accent in cards:
        border = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(8.8), cy, Inches(0.08), Inches(1.5))
        border.fill.solid(); border.fill.fore_color.rgb = accent; border.line.fill.background()
        text(s, Inches(9.0), cy + Inches(0.05), Inches(4), Inches(0.6),
             big, size=28, bold=True, color=accent)
        text(s, Inches(9.0), cy + Inches(0.65), Inches(4), Inches(0.8),
             lbl, size=11, color=DIM)
        cy += Inches(1.7)
    pagenum(s, 6)


# ============================================================
# 7 — OOD + REALTIME
# ============================================================
def s7(prs):
    s = new_slide(prs)
    title_bar(s, "When it fails  +  what realtime would buy")
    picture(s, Inches(0.4), Inches(1.6), Inches(6.2), Inches(3.6),
            "results/analysis/figures/07_in_vs_ood_delta.png")
    text(s, Inches(0.4), Inches(5.3), Inches(6.2), Inches(1.6),
         "OOD kill. Warehouse-trained model on residential maps: advantage drops from +3.5pp to −0.1pp. "
         "The prior is doing real structural work, not averaging.",
         size=11, italic=True, color=DIM)
    picture(s, Inches(7.0), Inches(1.6), Inches(5.9), Inches(3.6),
            "results/analysis/figures/11_2x2_matrix.png")
    text(s, Inches(7.0), Inches(5.3), Inches(5.9), Inches(1.6),
         "Realtime upper bound. K=4 / DDIM=10 = ~500 ms on T4. "
         "Sampling every 3 cells lifts advantage to +10.0pp. Distillation is the path.",
         size=11, italic=True, color=DIM)
    pagenum(s, 7)


# ============================================================
# 8 — HARDSHIPS
# ============================================================
def s8(prs):
    s = new_slide(prs)
    title_bar(s, "Hardships  +  Honest limits")
    items = [
        "**Docker QEMU on Apple Silicon** too slow for live diffusion in Stage — pre-computed waypoints instead.",
        "**Asymptotic catch-up** — baseline catches up by step 20.",
        "**2 / 30 hard failures** — K-sample disagreement could flag these.",
        "**Live hardware deployment** blocked by latency on Apple Silicon emulation.",
        "**Mismatched priors** give zero help — the OOD ablation confirms this.",
    ]
    bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(5), items, size=17, bullet_color=RED)
    pagenum(s, 8)


# ============================================================
# 9 — WHAT'S NEXT
# ============================================================
def s9(prs):
    s = new_slide(prs)
    title_bar(s, "What's Next")
    items = [
        "**Distillation** of the U-Net to Jetson for sub-second inference on a real robot.",
        "**Richer training domains** — Gibson, Matterport, BIM — for cross-environment generalisation.",
        "**K-sample disagreement** as an \"I don't know\" flag — fall back to heuristic on hard maps.",
        "**Live ROSbot deployment** in a matched-distribution building.",
    ]
    bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(5), items, size=18)
    pagenum(s, 9)


# ============================================================
# 10 — THANKS
# ============================================================
def s10(prs):
    s = new_slide(prs)
    text(s, Inches(0.6), Inches(2.9), Inches(12), Inches(1.2),
         "Thanks for your attention.", size=54, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    text(s, Inches(0.6), Inches(4.1), Inches(12), Inches(0.8),
         "Questions?", size=36, italic=True, color=GREEN, align=PP_ALIGN.CENTER)
    text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
         "Moin Mattar  ·  COSC 81/281  ·  June 1, 2026",
         size=14, color=DIM, align=PP_ALIGN.CENTER)
    text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
         "AI helped me in formatting and writing (HTML/CSS/LaTeX), as well as explained concepts.",
         size=10, color=DIM, align=PP_ALIGN.CENTER)
    pagenum(s, 10)


def main():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    for fn in [s1, s2, s3, s4, s5, s6, s7, s8, s9, s10]:
        fn(prs)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
