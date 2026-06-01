"""Two-slide PPTX matching the user's earlier presentation style:
white background, black text, green accents, boxed hypothesis callout.

Output: intro_slides.pptx (slide 1: Project Goal, slide 2: Novelty & Hypothesis)
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path("/Users/moin/Robots_diffusion_planner")
OUT = ROOT / "intro_slides.pptx"

W = Inches(13.333)
H = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DIM = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x1F, 0x6B, 0x47)
ACCENT = RGBColor(0x2D, 0x8A, 0x5C)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    return slide


def add_title(slide, text, x=Inches(0.7), y=Inches(0.55), color=BLACK, size=44):
    tb = slide.shapes.add_textbox(x, y, Inches(12), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = "Calibri"
    # underline bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.85), Inches(12), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()


def add_subtitle(slide, text, x=Inches(0.7), y=Inches(1.55), color=DIM, size=18):
    tb = slide.shapes.add_textbox(x, y, Inches(12), Inches(0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.italic = True
    r.font.color.rgb = color
    r.font.name = "Calibri"


def add_bullets(slide, x, y, w, h, items, size=18, color=BLACK, bullet_char="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(14)

        # parse simple bold marker: text with **bold** segments
        parts = []
        rest = item
        while "**" in rest:
            before, _, rest2 = rest.partition("**")
            bold_text, _, rest = rest2.partition("**")
            if before:
                parts.append((before, False))
            parts.append((bold_text, True))
        if rest:
            parts.append((rest, False))

        # bullet
        br = p.add_run()
        br.text = f"{bullet_char}  "
        br.font.size = Pt(size)
        br.font.color.rgb = GREEN
        br.font.bold = True
        br.font.name = "Calibri"

        for text, bold in parts:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"


def add_pagenum(slide, n):
    tb = slide.shapes.add_textbox(W - Inches(0.7), H - Inches(0.45), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(n)
    r.font.size = Pt(11)
    r.font.color.rgb = DIM
    r.font.name = "Calibri"


def add_image(slide, x, y, w, h, path):
    p = ROOT / path
    if not p.exists():
        print(f"Warning: missing {p}")
        return None
    return slide.shapes.add_picture(str(p), x, y, w, h)


def slide_1(prs):
    s = new_slide(prs)
    add_title(s, "Project Goal")
    add_subtitle(s, "Autonomous mapping of an unknown building, one frontier decision at a time.")

    bullets = [
        "**Setup:** a robot is dropped at the door of a building it has never seen, with only a 2D lidar.",
        "**Input each step:** the partial occupancy grid the robot has accumulated so far.",
        "**Output each step:** the next frontier to drive to — the boundary between known and unknown.",
        "**Assumption:** the deployment environment shares structural distribution with training data (real residential floor plans).",
        "**What's broken with the standard approach:** classical frontier scoring is memoryless — it picks based on nearby unknown cells minus distance, throwing away the fact that real buildings have walls that run straight, hallways that connect rooms, doors that cluster.",
    ]
    add_bullets(s, Inches(0.7), Inches(2.3), Inches(7.8), Inches(5), bullets, size=15)

    # Right-side image: behind-mind composite (cropped to show structure)
    add_image(s, Inches(8.9), Inches(2.3), Inches(3.9), Inches(4.4),
              "presentation_media/figures/slide11_behind_mind_map2638.png")

    # Caption under image
    tb = s.shapes.add_textbox(Inches(8.9), Inches(6.75), Inches(3.9), Inches(0.4))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "the structural prior the heuristic ignores"
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = DIM
    r.font.name = "Calibri"

    add_pagenum(s, 1)


def slide_2(prs):
    s = new_slide(prs)
    add_title(s, "Novelty & Hypothesis")

    # Three novelty bullets (top half)
    novelty_label = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8), Inches(0.4))
    tf = novelty_label.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "What's new here"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = GREEN
    r.font.name = "Calibri"

    bullets = [
        "We treat **diffusion samples** as imagined buildings, not just predictions. K=8 plausible completions per frontier decision.",
        "We use **K-sample disagreement** as an upper-confidence-bound exploration signal — the variance is the feature, not noise to be averaged out.",
        "We isolate **where the wins come from** with a 4-baseline ablation (nearest / info-gain / info-gain + distance / ours), plus an out-of-domain control that proves the prior is doing real structural work.",
    ]
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(8.1), Inches(2.6), bullets, size=14)

    # Right side small image: K diversity
    add_image(s, Inches(9.2), Inches(1.8), Inches(3.7), Inches(2.7),
              "presentation_media/figures/slide07_K_diversity.png")

    # Hypothesis label
    hyp_label = s.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(8), Inches(0.4))
    tf = hyp_label.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Hypothesis"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = BLACK
    r.font.name = "Calibri"

    # Boxed hypothesis with green left border (matching earlier deck style)
    border = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(5.5), Inches(0.08), Inches(1.6))
    border.fill.solid(); border.fill.fore_color.rgb = GREEN
    border.line.fill.background()

    box = s.shapes.add_textbox(Inches(0.95), Inches(5.5), Inches(12.0), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = 0; tf.margin_top = Inches(0.05); tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.35

    r1 = p.add_run()
    r1.text = "A learned generative prior of building layouts gives a robot a measurable early-budget head start in frontier exploration "
    r1.font.size = Pt(17); r1.font.color.rgb = BLACK; r1.font.name = "Calibri"

    r2 = p.add_run()
    r2.text = "— but only when "
    r2.font.size = Pt(17); r2.font.italic = True; r2.font.color.rgb = BLACK; r2.font.name = "Calibri"

    r3 = p.add_run()
    r3.text = "(1) the training distribution matches deployment, (2) the step budget is tight, and (3) inference is fast enough to drive."
    r3.font.size = Pt(17); r3.font.bold = True; r3.font.color.rgb = BLACK; r3.font.name = "Calibri"

    add_pagenum(s, 2)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_1(prs)
    slide_2(prs)

    prs.save(str(OUT))
    print(f"Wrote {OUT} ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
